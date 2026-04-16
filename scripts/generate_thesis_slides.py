#!/usr/bin/env python3
"""
=============================================================================
Generate Thesis Defense Slides (.pptx)
=============================================================================

Generates a 13-slide PowerPoint deck (11 content + 2 backup) summarizing
the multi-agent LLM negotiation research across Games 1, 2, and 3.

Usage:
    python scripts/generate_thesis_slides.py

What it creates:
    slides/thesis_defense.pptx

Dependencies:
    python-pptx, Pillow (via python-pptx)

=============================================================================
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── paths ──
ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "Figures"
VIZ = ROOT / "visualization" / "figures"
N2_ANALYSIS = (
    ROOT
    / "experiments"
    / "results"
    / "game1_multiagent_full_20260413_045538"
    / "analysis"
)
CROSS_GAME = VIZ / "gpt5_nano_baseline_vs_elo_all_games_20260413"
OUT_DIR = ROOT / "slides"
OUT_DIR.mkdir(exist_ok=True)

# ── colours ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
GRAY = RGBColor(0x6C, 0x75, 0x7D)
DARK_GRAY = RGBColor(0x34, 0x3A, 0x40)
LIGHT_GRAY = RGBColor(0xDE, 0xE2, 0xE6)
GREEN = RGBColor(0x28, 0xA7, 0x45)
RED = RGBColor(0xDC, 0x35, 0x45)
ORANGE = RGBColor(0xFD, 0x7E, 0x14)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background colour for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri", line_spacing=1.15):
    """Add a text box to a slide with formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return tf


def add_bullet_frame(slide, left, top, width, height, items, font_size=16,
                     color=BLACK, font_name="Calibri", bold_items=None,
                     line_spacing=1.3):
    """Add a bulleted text frame. bold_items is a set of indices to bold."""
    if bold_items is None:
        bold_items = set()
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.font.bold = i in bold_items
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        p.line_spacing = Pt(font_size * line_spacing)
        # bullet
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
        # remove existing bullets first
        for old in pPr.findall(qn("a:buChar")):
            pPr.remove(old)
        for old in pPr.findall(qn("a:buNone")):
            pPr.remove(old)
        pPr.append(buChar)
    return tf


def add_image_centered(slide, img_path, top, max_width, max_height):
    """Add an image centered horizontally, scaled to fit within bounds."""
    if not os.path.exists(img_path):
        add_textbox(slide, Inches(2), top, max_width, Inches(1),
                    f"[Figure not found: {os.path.basename(img_path)}]",
                    font_size=14, color=RED, alignment=PP_ALIGN.CENTER)
        return
    from PIL import Image
    with Image.open(img_path) as img:
        iw, ih = img.size
    aspect = iw / ih
    # fit within bounds
    w = max_width
    h = w / aspect
    if h > max_height:
        h = max_height
        w = h * aspect
    left = (SLIDE_WIDTH - w) / 2
    slide.shapes.add_picture(str(img_path), int(left), int(top), int(w), int(h))


def add_divider(slide, top, color=ACCENT_BLUE, width=None):
    """Add a horizontal line."""
    if width is None:
        width = SLIDE_WIDTH - Inches(2)
    left = (SLIDE_WIDTH - width) / 2
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, int(left), int(top), int(width), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_slide_number(slide, num, total=13):
    """Add slide number at bottom right."""
    add_textbox(slide, SLIDE_WIDTH - Inches(1.5), SLIDE_HEIGHT - Inches(0.5),
                Inches(1.2), Inches(0.4), f"{num}/{total}",
                font_size=11, color=GRAY, alignment=PP_ALIGN.RIGHT)


# ======================================================================
#  BUILD DECK
# ======================================================================

prs = Presentation()
prs.slide_width = int(SLIDE_WIDTH)
prs.slide_height = int(SLIDE_HEIGHT)

blank_layout = prs.slide_layouts[6]  # blank

# ----------------------------------------------------------------------
# SLIDE 1: Title (0.5 min)
# ----------------------------------------------------------------------
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, DARK_BLUE)

add_textbox(sl, Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
            "Scaling Laws for LLM Negotiation",
            font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
            font_name="Calibri")

add_divider(sl, Inches(3.1), color=ACCENT_BLUE, width=Inches(6))

add_textbox(sl, Inches(1), Inches(3.5), Inches(11.3), Inches(1.2),
            "How model capability shapes bargaining outcomes\nacross cooperative and competitive regimes",
            font_size=22, color=RGBColor(0xAE, 0xBF, 0xD5), alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1), Inches(5.2), Inches(11.3), Inches(0.5),
            "Jonathan Zhu",
            font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1), Inches(5.7), Inches(11.3), Inches(0.5),
            "Advisor: Prof. Danqi Chen  |  Princeton University  |  2026",
            font_size=16, color=RGBColor(0x8E, 0x9F, 0xB5), alignment=PP_ALIGN.CENTER)

# Claim
add_textbox(sl, Inches(2), Inches(6.4), Inches(9.3), Inches(0.8),
            "Thesis: Higher Elo reliably predicts better negotiation payoffs, "
            "but competition structure and game mechanism modulate the slope.",
            font_size=15, color=RGBColor(0xCE, 0xD9, 0xE5), alignment=PP_ALIGN.CENTER)

# ----------------------------------------------------------------------
# SLIDE 2: Motivation (1.5 min)
# ----------------------------------------------------------------------
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, WHITE)

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
            "Why Study LLM Negotiation?",
            font_size=36, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
add_divider(sl, Inches(1.15), color=ACCENT_BLUE, width=Inches(11.7))

bullets = [
    "LLMs are increasingly deployed as autonomous agents in multi-party settings",
    "Negotiation is a canonical testbed for strategic reasoning in natural language",
    "Prior work shows stronger models diverge into model-specific strategic signatures "
    "rather than converging to a common policy (Rios et al., 2025)",
    "Scaling laws well-studied for single-agent RL, vision, and LLMs \u2014 "
    "but not for multi-agent negotiation with tunable strategic structure",
    "Key safety question: does higher capability lead to exploitation of weaker partners?",
]
add_bullet_frame(sl, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.5),
                 bullets, font_size=20, color=DARK_GRAY, bold_items={4})

# Bottom accent box
shape = sl.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.7), Inches(10.3), Inches(1.2)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Gap: No systematic study of how LLM capability interacts with "
          "cooperation\u2013competition structure across multiple negotiation mechanisms")
p.font.size = Pt(17)
p.font.color.rgb = ACCENT_BLUE
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

add_slide_number(sl, 2)

# ----------------------------------------------------------------------
# SLIDE 3: Research Questions & Design (1.5 min)
# ----------------------------------------------------------------------
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, WHITE)

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
            "Research Questions & Experiment Design",
            font_size=36, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
add_divider(sl, Inches(1.15), color=ACCENT_BLUE, width=Inches(11.7))

# Research questions
rqs = [
    "RQ1: Does higher Elo translate into better bargaining outcomes against a fixed baseline?",
    "RQ2: How much variation comes from competition parameters vs. capability alone?",
    "RQ3: What changes across three distinct negotiation mechanisms?",
]
add_bullet_frame(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5),
                 rqs, font_size=18, color=DARK_GRAY, bold_items={0, 1, 2})

# Design summary box
design = [
    "Fixed baseline: GPT-5-nano (Elo 1338)",
    "36 adversary models spanning Elo range \u2212226 to +152",
    "3 game mechanisms with tunable competition",
    "n = 2 agents, m = 5 items/issues, T = 10 rounds",
    "~1,970 total completed game configurations",
]
add_bullet_frame(sl, Inches(6.8), Inches(1.5), Inches(5.7), Inches(3.0),
                 design, font_size=17, color=DARK_GRAY)

# Three game boxes
games_info = [
    ("Game 1: Item Allocation", "Rivalrous goods\n\u03b1 \u2208 [0,1] competition", ACCENT_BLUE),
    ("Game 2: Diplomatic Treaty", "Continuous terms\n\u03c1, \u03b8 competition axes", GREEN),
    ("Game 3: Co-Funding", "Threshold public goods\n\u03b1, \u03c3 parameters", ORANGE),
]
box_width = Inches(3.5)
box_height = Inches(2.0)
box_top = Inches(4.8)
gap = (Inches(11.7) - 3 * box_width) / 2

for i, (title, desc, accent) in enumerate(games_info):
    left = Inches(0.8) + i * (box_width + gap)
    shape = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                int(left), int(box_top),
                                int(box_width), int(box_height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = accent
    shape.line.width = Pt(2)

    # Title
    add_textbox(sl, left + Inches(0.2), box_top + Inches(0.15),
                box_width - Inches(0.4), Inches(0.5),
                title, font_size=16, bold=True, color=accent,
                alignment=PP_ALIGN.CENTER)
    # Description
    add_textbox(sl, left + Inches(0.2), box_top + Inches(0.7),
                box_width - Inches(0.4), Inches(1.0),
                desc, font_size=14, color=DARK_GRAY,
                alignment=PP_ALIGN.CENTER)

add_slide_number(sl, 3)

# ----------------------------------------------------------------------
# SLIDE 4: Game 1 Environment (2 min)
# ----------------------------------------------------------------------
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, WHITE)

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
            "Game 1: Item Allocation",
            font_size=36, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
add_divider(sl, Inches(1.15), color=ACCENT_BLUE, width=Inches(11.7))

# Left column - math description
math_bullets = [
    "n agents split m items; each agent has private valuations v_i (sum to 100)",
    "Utility: u_i(A) = \u03b3^(t-1) \u00b7 A[i] \u00b7 v_i  (discount \u03b3 = 0.9)",
    "Competition: cosine similarity \u03b1 = cos(v_i, v_j) \u2208 [0,1]",
    "\u03b1 = 1 \u2192 zero-sum (same preferences); \u03b1 = 0 \u2192 cooperative (orthogonal)",
    "Propose-and-vote protocol: discuss \u2192 think \u2192 propose \u2192 vote \u2192 reflect",
    "Unanimous approval required; T = 10 rounds, then utility = 0",
]
add_bullet_frame(sl, Inches(0.8), Inches(1.4), Inches(5.8), Inches(3.5),
                 math_bullets, font_size=16, color=DARK_GRAY)

# Right column - figure
add_image_centered(sl, FIG / "game_1" / "average_utility_vs_elo.png",
                   Inches(1.3), Inches(6.0), Inches(5.5))

add_textbox(sl, Inches(6.5), Inches(6.7), Inches(6.0), Inches(0.5),
            "890 completed runs  |  30 adversary models  |  r = 0.85",
            font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(sl, 4)

# ----------------------------------------------------------------------
# SLIDE 5: Game 2 Environment (2 min)
# ----------------------------------------------------------------------
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, WHITE)

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
            "Game 2: Diplomatic Treaty Negotiation",
            font_size=36, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
add_divider(sl, Inches(1.15), color=ACCENT_BLUE, width=Inches(11.7))

math_bullets = [
    "n agents negotiate m continuous issues; agreement a \u2208 [0,1]^m",
    "Each agent has position prefs p_i and importance weights w_i",
    "Utility: u_i(a) = 100 \u00b7 \u03a3 w_ik \u00b7 (1 - |p_ik - a_k|)",
    "Two competition axes (Gaussian copula preferences):",
    "  \u03c1 \u2208 [\u22120.25, 1]: preference correlation (same vs. opposing outcomes)",
    "  \u03b8 \u2208 [0, 1]: interest overlap (same vs. different issues)",
    "Same propose-and-vote protocol as Game 1",
]
add_bullet_frame(sl, Inches(0.8), Inches(1.4), Inches(5.8), Inches(3.5),
                 math_bullets, font_size=16, color=DARK_GRAY)

# Figure
add_image_centered(sl, FIG / "game_2" / "utility_vs_elo_by_rho_theta.png",
                   Inches(1.3), Inches(6.0), Inches(5.5))

add_textbox(sl, Inches(6.5), Inches(6.7), Inches(6.0), Inches(0.5),
            "540/540 completed  |  30 models  |  r = 0.82",
            font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(sl, 5)

# ----------------------------------------------------------------------
# SLIDE 6: Game 3 Environment (2.5 min)
# ----------------------------------------------------------------------
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, WHITE)

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
            "Game 3: Participatory Budgeting (Co-Funding)",
            font_size=36, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
add_divider(sl, Inches(1.15), color=ACCENT_BLUE, width=Inches(11.7))

math_bullets = [
    "n agents pool budgets to fund m projects (threshold public goods)",
    "Project j funded iff total contributions \u2265 cost c_j (refund if not)",
    "Utility: u_i(X) = \u03a3_{j\u2208S} v_ij - \u03a3_{j\u2208S} x_ij  (value minus cost share)",
    "Key difference: threshold non-linearity creates free-riding incentives",
    "Two competition axes:",
    "  \u03b1 \u2208 [0,1]: preference alignment (shared vs. disjoint project values)",
    "  \u03c3 \u2208 (0,1]: budget scarcity (budget / total cost)",
    "Talk\u2013Pledge\u2013Revise protocol (not propose-and-vote)",
]
add_bullet_frame(sl, Inches(0.8), Inches(1.4), Inches(5.8), Inches(4.0),
                 math_bullets, font_size=15, color=DARK_GRAY)

# Figure
add_image_centered(sl, FIG / "game_3" / "utility_vs_elo_by_alpha_sigma.png",
                   Inches(1.3), Inches(6.0), Inches(5.5))

add_textbox(sl, Inches(6.5), Inches(6.7), Inches(6.0), Inches(0.5),
            "540/540 completed  |  29 adversary models  |  r = 0.72",
            font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(sl, 6)

# ----------------------------------------------------------------------
# SLIDE 7: Evaluation Metrics & Benchmarks (1.5 min)
# ----------------------------------------------------------------------
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, WHITE)

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
            "Evaluation Metrics & Solution Benchmarks",
            font_size=36, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
add_divider(sl, Inches(1.15), color=ACCENT_BLUE, width=Inches(11.7))

# Left: metrics
add_textbox(sl, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
            "Core Metrics", font_size=22, bold=True, color=ACCENT_BLUE)
metrics = [
    "Utilitarian efficiency \u03b7 = SW_actual / SW_optimal",
    "Exploitation index E_i = (u_actual - u_fair) / |u_fair|",
    "Efficiency\u2013fairness decomposition:\n  - Efficiency loss = SW* - SW_actual\n  - Fairness deviation = distance from benchmark",
]
add_bullet_frame(sl, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.0),
                 metrics, font_size=16, color=DARK_GRAY)

# Right: benchmarks
add_textbox(sl, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.5),
            "Fair Benchmarks", font_size=22, bold=True, color=ACCENT_BLUE)
benchmarks = [
    "Games 1 & 2: Nash Bargaining Solution (NBS)\n  max \u03a0 u_i(A) \u2014 unique, Pareto efficient",
    "Game 3: Lindahl equilibrium\n  x_ij = c_j \u00b7 v_ij / \u03a3_k v_kj \u2014 proportional cost-sharing",
    "Game 3: Core of cooperative game\n  No coalition can do better independently",
    "Optimal computed exactly: greedy (G1), weighted median (G2), knapsack (G3)",
]
add_bullet_frame(sl, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.0),
                 benchmarks, font_size=15, color=DARK_GRAY)

# Bottom summary
add_textbox(sl, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8),
            "Normalization: All games on [0, 100] utility scale. "
            "Cross-\u03c3 comparisons use \u03b7. Exploitation index E_i comparable across games.",
            font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(sl, 7)

# ----------------------------------------------------------------------
# SLIDE 8: Game 1 Results (2 min)
# ----------------------------------------------------------------------
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, WHITE)

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
            "Game 1 Results: Capability Buys Better Deals, Faster",
            font_size=34, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
add_divider(sl, Inches(1.15), color=ACCENT_BLUE, width=Inches(11.7))

# Two figures side by side
fig_top = Inches(1.4)
fig_width = Inches(5.8)
fig_height = Inches(4.5)

# Left figure: utility vs competition
left_fig = FIG / "game_1" / "average_utility_vs_competition_level_aggregated_over_models.png"
if os.path.exists(left_fig):
    from PIL import Image
    with Image.open(left_fig) as img:
        iw, ih = img.size
    aspect = iw / ih
    w = fig_width
    h = int(int(w) / aspect)
    if h > int(fig_height):
        h = int(fig_height)
        w = int(h * aspect)
    sl.shapes.add_picture(str(left_fig), Inches(0.5), int(fig_top), int(w), int(h))

# Right figure: rounds vs elo
right_fig = FIG / "game_1" / "average_rounds_to_consensus_vs_elo.png"
if os.path.exists(right_fig):
    from PIL import Image
    with Image.open(right_fig) as img:
        iw, ih = img.size
    aspect = iw / ih
    w = fig_width
    h = int(int(w) / aspect)
    if h > int(fig_height):
        h = int(fig_height)
        w = int(h * aspect)
    sl.shapes.add_picture(str(right_fig), Inches(6.8), int(fig_top), int(w), int(h))

# Key findings
findings = [
    "Utility drops from 95.1 (\u03b1=0) to 39.4 (\u03b1=1)  |  "
    "Strongest models preserve leverage even at \u03b1=1  |  "
    "Elo\u2013rounds correlation: \u22120.47 (faster deals)",
]
add_bullet_frame(sl, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8),
                 findings, font_size=15, color=DARK_GRAY)

add_slide_number(sl, 8)

# ----------------------------------------------------------------------
# SLIDE 9: Game 2 & Game 3 Results (2 min)
# ----------------------------------------------------------------------
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, WHITE)

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
            "Game 2 & Game 3 Results",
            font_size=36, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
add_divider(sl, Inches(1.15), color=ACCENT_BLUE, width=Inches(11.7))

# Two columns
add_textbox(sl, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
            "Game 2: Treaty (540/540 complete)", font_size=20, bold=True, color=GREEN)

g2_bullets = [
    "r = 0.82 Elo\u2013utility correlation",
    "Much higher utility level: mean ~85 across all conditions",
    "\u03c1 (preference correlation) matters more than \u03b8 (interest overlap)",
    "\u03c1=1 \u2192 utility 97.6; \u03c1=\u22121 \u2192 utility 74.7",
    "Most forgiving: even hardest condition, mean utility > 65",
    "All runs reach agreement; mean 1.5 rounds",
]
add_bullet_frame(sl, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.0),
                 g2_bullets, font_size=15, color=DARK_GRAY)

add_textbox(sl, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.5),
            "Game 3: Co-Funding (540/540 complete)", font_size=20, bold=True, color=ORANGE)

g3_bullets = [
    "r = 0.72 Elo\u2013utility (weakest of three games)",
    "Much lower payoffs: GPT-5.4 High leads at only 40.5",
    "\u03c3 (scarcity) does most of the damage, not \u03b1 (alignment)",
    "Hardest cell: some models go negative (lost money!)",
    "Threshold non-linearity \u2192 coordination failure risk",
    "Nova Micro averages 0.26 utility; Llama 3.1 8B goes to \u22124.46",
]
add_bullet_frame(sl, Inches(7.0), Inches(2.0), Inches(5.5), Inches(3.0),
                 g3_bullets, font_size=15, color=DARK_GRAY)

# Bottom comparison box
shape = sl.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.5), Inches(11.3), Inches(1.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFD, 0xF3, 0xE7)
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Key contrast: Treaty bargaining is forgiving (continuous compromise always possible). "
          "Co-funding is harsh (threshold non-linearity means small mistakes \u2192 zero payoff). "
          "Game mechanism determines how severe competition's impact becomes.")
p.font.size = Pt(16)
p.font.color.rgb = DARK_GRAY
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

add_slide_number(sl, 9)

# ----------------------------------------------------------------------
# SLIDE 10: Cross-Game Synthesis (1.5 min)
# ----------------------------------------------------------------------
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, WHITE)

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
            "Cross-Game Synthesis",
            font_size=36, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
add_divider(sl, Inches(1.15), color=ACCENT_BLUE, width=Inches(11.7))

# Three-column comparison
col_data = [
    ("Game 1\nItem Allocation", ACCENT_BLUE, [
        "r = 0.85",
        "Steady decline with competition",
        "Stronger models: better splits + faster deals",
        "Top: Opus 4.5 (72.6)",
    ]),
    ("Game 2\nDiplomatic Treaty", GREEN, [
        "r = 0.82",
        "Milder competition effect",
        "High utility even in hardest conditions",
        "Top: Opus 4.6T (94.4)",
    ]),
    ("Game 3\nCo-Funding", ORANGE, [
        "r = 0.72",
        "Sharpest competition effect",
        "Coordination failure at extremes",
        "Top: GPT-5.4H (40.5)",
    ]),
]

col_width = Inches(3.5)
col_top = Inches(1.5)

for i, (title, accent, items) in enumerate(col_data):
    left = Inches(0.8) + i * Inches(4.1)

    add_textbox(sl, left, col_top, col_width, Inches(0.8),
                title, font_size=18, bold=True, color=accent,
                alignment=PP_ALIGN.CENTER)

    add_bullet_frame(sl, left, col_top + Inches(0.9), col_width, Inches(2.5),
                     items, font_size=15, color=DARK_GRAY)

# Cross-game figure if available
cross_fig = CROSS_GAME / "game1_gpt5_nano_baseline_payoff_vs_adversary_elo_overall.png"
if os.path.exists(cross_fig):
    add_image_centered(sl, cross_fig, Inches(4.3), Inches(7.0), Inches(2.8))
    add_textbox(sl, Inches(2), Inches(6.9), Inches(9.3), Inches(0.3),
                "GPT-5-nano baseline payoff decreases as adversary Elo rises \u2014 "
                "consistent across all three games",
                font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
else:
    # Synthesis text instead
    shape = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.8), Inches(11.3), Inches(2.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ("Consistent finding across all three games:\n"
              "1. Higher Elo \u2192 better negotiation payoffs (always)\n"
              "2. Competition lowers utility for everyone (mechanism determines severity)\n"
              "3. Stronger models use their advantage differently per mechanism:\n"
              "   - G1: better splits + faster consensus\n"
              "   - G2: better terms (already fast convergence)\n"
              "   - G3: better coordination under scarcity (key differentiator)")
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_GRAY
    p.font.name = "Calibri"

add_slide_number(sl, 10)

# ----------------------------------------------------------------------
# SLIDE 11: Limitations & Future Work (1 min)
# ----------------------------------------------------------------------
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, WHITE)

add_textbox(sl, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
            "Limitations & Future Work",
            font_size=36, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
add_divider(sl, Inches(1.15), color=ACCENT_BLUE, width=Inches(11.7))

# Limitations
add_textbox(sl, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
            "Current Limitations", font_size=22, bold=True, color=RED)
lim_bullets = [
    "All results are n = 2 (bilateral negotiation only)",
    "Fixed baseline (GPT-5-nano) \u2014 results may differ with other baselines",
    "Confounds: instruction-following, sycophancy vs. true strategic capability",
    "API-based models: no control over decoding, context length varies",
    "Single-shot competition parameters per game instance",
]
add_bullet_frame(sl, Inches(0.8), Inches(2.0), Inches(5.5), Inches(3.5),
                 lim_bullets, font_size=16, color=DARK_GRAY)

# Future work
add_textbox(sl, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.5),
            "Future Work", font_size=22, bold=True, color=GREEN)
future_bullets = [
    "n > 2 extension: 9,600 configs currently running on cluster\n"
    "  (316 completed as of April 14, 2026)",
    "Multi-domain expansion beyond negotiation (ICML target)",
    "Behavioral metrics pipeline: promise-keeping, free-riding, persuasion",
    "Adaptive opponents (profit-driven red teaming)",
    "Test-time compute scaling for strategic reasoning",
]
add_bullet_frame(sl, Inches(7.0), Inches(2.0), Inches(5.5), Inches(3.5),
                 future_bullets, font_size=16, color=DARK_GRAY)

# Bottom
add_textbox(sl, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7),
            "Thank you! \u2014 Questions?",
            font_size=28, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

add_slide_number(sl, 11)

# ======================================================================
#  BACKUP SLIDES
# ======================================================================

# ----------------------------------------------------------------------
# BACKUP A: Game 2 Heatmaps
# ----------------------------------------------------------------------
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, LIGHT_BG)

add_textbox(sl, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.8),
            "Backup A: Game 2 \u2014 \u03c1\u00d7\u03b8 Heatmaps",
            font_size=30, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
add_divider(sl, Inches(0.95), color=GRAY, width=Inches(11.7))

heatmap_dir = VIZ / "diplomacy_20260405_082215_summary"
h1 = heatmap_dir / "utility_rho_theta_heatmap.png"
h2 = heatmap_dir / "social_welfare_rho_theta_heatmap.png"

fig_top = Inches(1.2)
fig_w = Inches(5.8)
fig_h = Inches(5.0)

if os.path.exists(h1):
    from PIL import Image
    with Image.open(h1) as img:
        iw, ih = img.size
    aspect = iw / ih
    w = fig_w
    h = int(int(w) / aspect)
    if h > int(fig_h):
        h = int(fig_h)
        w = int(h * aspect)
    sl.shapes.add_picture(str(h1), Inches(0.3), int(fig_top), int(w), int(h))

if os.path.exists(h2):
    from PIL import Image
    with Image.open(h2) as img:
        iw, ih = img.size
    aspect = iw / ih
    w = fig_w
    h = int(int(w) / aspect)
    if h > int(fig_h):
        h = int(fig_h)
        w = int(h * aspect)
    sl.shapes.add_picture(str(h2), Inches(6.8), int(fig_top), int(w), int(h))

add_textbox(sl, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5),
            "Left: Mean adversary utility  |  Right: Social welfare  |  "
            "Both show \u03c1 drives more variation than \u03b8",
            font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(sl, "B-A", total="B")

# ----------------------------------------------------------------------
# BACKUP B: n > 2 Status
# ----------------------------------------------------------------------
sl = prs.slides.add_slide(blank_layout)
set_slide_bg(sl, LIGHT_BG)

add_textbox(sl, Inches(0.8), Inches(0.3), Inches(11.7), Inches(0.8),
            "Backup B: n > 2 Extension \u2014 Early Results",
            font_size=30, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT)
add_divider(sl, Inches(0.95), color=GRAY, width=Inches(11.7))

status_bullets = [
    "Game 1 multi-agent batch: 9,600 total configurations",
    "As of April 14, 2026: 316 successful, 181 running, 1,506 queued, 7,596 not started",
    "Early finding: Proposal-1 advantage varies with Elo (see figure)",
    "Full results expected before camera-ready deadline",
]
add_bullet_frame(sl, Inches(0.8), Inches(1.2), Inches(5.5), Inches(2.5),
                 status_bullets, font_size=17, color=DARK_GRAY)

# Proposal 1 advantage figure
prop1_fig = N2_ANALYSIS / "proposal1_advantage_vs_elo.png"
if os.path.exists(prop1_fig):
    add_image_centered(sl, prop1_fig, Inches(1.2), Inches(6.5), Inches(5.5))

add_textbox(sl, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
            "Preliminary: batch incomplete \u2014 treat as ongoing extension, not finalized result",
            font_size=14, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)

add_slide_number(sl, "B-B", total="B")

# ======================================================================
#  SAVE
# ======================================================================
out_path = OUT_DIR / "thesis_defense.pptx"
prs.save(str(out_path))
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)} (11 content + 2 backup)")
